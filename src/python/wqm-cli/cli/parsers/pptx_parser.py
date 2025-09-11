from ...observability import get_logger

logger = get_logger(__name__)
"""
Microsoft PowerPoint PPTX presentation parser for extracting text content and metadata.

This module provides functionality to parse PPTX PowerPoint presentations, extracting
text content from slides including titles, bullet points, speaker notes, and embedded
text while preserving presentation metadata and structure information.
"""

import logging
from pathlib import Path
from typing import TYPE_CHECKING, Any

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.shapes.autoshape import Shape
    from pptx.shapes.placeholder import PlaceholderPicture, SlidePlaceholder
    from pptx.slide import Slide

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    if TYPE_CHECKING:
        from pptx import Presentation
        from pptx.slide import Slide

from .base import DocumentParser, ParsedDocument

logger = logging.getLogger(__name__)


class PptxParser(DocumentParser):
    """
    Parser for Microsoft PowerPoint PPTX presentations.

    Extracts text content from slides including titles, bullet points, and speaker notes
    while preserving presentation metadata like title, author, and creation date.
    """

    @property
    def supported_extensions(self) -> list[str]:
        """PPTX file extensions."""
        return [".pptx"]

    @property
    def format_name(self) -> str:
        """Human-readable format name."""
        return "Microsoft PowerPoint PPTX"

    def _check_availability(self) -> None:
        """Check if required libraries are available."""
        if not PPTX_AVAILABLE:
            raise RuntimeError(
                "PPTX parsing requires 'python-pptx'. "
                "Install with: pip install python-pptx"
            )

    async def parse(self, file_path: str | Path, **options: Any) -> ParsedDocument:
        """
        Parse PPTX file and extract text content.

        Args:
            file_path: Path to PPTX file
            **options: Parsing options
                - include_speaker_notes: bool = True - Include speaker notes
                - include_slide_numbers: bool = True - Include slide numbers in output
                - slide_separator: str = "\n\n---\n\n" - Separator between slides
                - extract_table_content: bool = True - Extract text from tables
                - extract_chart_titles: bool = True - Extract titles from charts
                - include_hidden_slides: bool = False - Include hidden slides

        Returns:
            ParsedDocument with extracted text and metadata

        Raises:
            RuntimeError: If required libraries are not installed
            ValueError: If file format is not supported
            FileNotFoundError: If file doesn't exist
        """
        self._check_availability()
        self.validate_file(file_path)

        file_path = Path(file_path)

        try:
            # Parse options
            include_speaker_notes = options.get("include_speaker_notes", True)
            include_slide_numbers = options.get("include_slide_numbers", True)
            slide_separator = options.get("slide_separator", "\n\n---\n\n")
            extract_table_content = options.get("extract_table_content", True)
            extract_chart_titles = options.get("extract_chart_titles", True)
            include_hidden_slides = options.get("include_hidden_slides", False)

            # Read PPTX presentation
            prs = Presentation(str(file_path))

            # Extract metadata
            metadata = await self._extract_metadata(prs)

            # Extract text content
            text_content = await self._extract_text_content(
                prs,
                include_speaker_notes,
                include_slide_numbers,
                slide_separator,
                extract_table_content,
                extract_chart_titles,
                include_hidden_slides,
            )

            # Count slides (including hidden if requested)
            visible_slides = sum(
                1
                for slide in prs.slides
                if not getattr(slide, "slide_id", None)
                or not slide.slide_layout.master.slide_master._element.get("visibility")
                == "0"
            )
            total_slides = len(prs.slides)
            processed_slides = total_slides if include_hidden_slides else visible_slides

            # Parsing information
            parsing_info = {
                "slide_count": total_slides,
                "visible_slide_count": visible_slides,
                "processed_slide_count": processed_slides,
                "slide_layout_count": len(prs.slide_layouts),
                "slide_master_count": len(prs.slide_masters),
                "content_length": len(text_content),
                "include_speaker_notes": include_speaker_notes,
                "include_slide_numbers": include_slide_numbers,
                "extract_table_content": extract_table_content,
                "extract_chart_titles": extract_chart_titles,
                "include_hidden_slides": include_hidden_slides,
            }

            logger.info(
                f"Successfully parsed PPTX: {file_path.name} "
                f"({parsing_info['processed_slide_count']} slides processed, "
                f"{parsing_info['slide_layout_count']} layouts, "
                f"{parsing_info['content_length']:,} characters)"
            )

            return ParsedDocument.create(
                content=text_content,
                file_path=file_path,
                file_type="pptx",
                additional_metadata=metadata,
                parsing_info=parsing_info,
            )

        except Exception as e:
            logger.error(f"Failed to parse PPTX {file_path}: {e}")
            raise RuntimeError(f"PPTX parsing failed: {e}") from e

    async def _extract_metadata(
        self, prs: "Presentation"
    ) -> dict[str, str | int | float | bool]:
        """Extract metadata from PPTX presentation."""
        metadata = {}

        # Core document properties
        core_props = prs.core_properties

        if core_props.title:
            metadata["title"] = core_props.title
        if core_props.author:
            metadata["author"] = core_props.author
        if core_props.subject:
            metadata["subject"] = core_props.subject
        if core_props.keywords:
            metadata["keywords"] = core_props.keywords
        if core_props.comments:
            metadata["comments"] = core_props.comments
        if core_props.category:
            metadata["category"] = core_props.category
        if core_props.language:
            metadata["language"] = core_props.language

        # Dates
        if core_props.created:
            metadata["created_date"] = core_props.created.isoformat()
        if core_props.modified:
            metadata["modified_date"] = core_props.modified.isoformat()
        if core_props.last_modified_by:
            metadata["last_modified_by"] = core_props.last_modified_by

        # Presentation statistics
        metadata["slide_count"] = len(prs.slides)
        metadata["slide_layout_count"] = len(prs.slide_layouts)
        metadata["slide_master_count"] = len(prs.slide_masters)

        # Slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        metadata["slide_width"] = slide_width
        metadata["slide_height"] = slide_height
        metadata["slide_aspect_ratio"] = (
            round(slide_width / slide_height, 2) if slide_height > 0 else 0
        )

        # Content analysis
        has_images = await self._has_images(prs)
        has_tables = await self._has_tables(prs)
        has_charts = await self._has_charts(prs)
        has_speaker_notes = await self._has_speaker_notes(prs)

        metadata["has_images"] = has_images
        metadata["has_tables"] = has_tables
        metadata["has_charts"] = has_charts
        metadata["has_speaker_notes"] = has_speaker_notes

        return metadata

    async def _extract_text_content(
        self,
        prs: "Presentation",
        include_speaker_notes: bool,
        include_slide_numbers: bool,
        slide_separator: str,
        extract_table_content: bool,
        extract_chart_titles: bool,
        include_hidden_slides: bool,
    ) -> str:
        """Extract text content from the presentation."""
        slide_contents = []

        for slide_idx, slide in enumerate(prs.slides, 1):
            # Skip hidden slides if not requested
            if not include_hidden_slides and self._is_slide_hidden(slide):
                continue

            slide_text_parts = []

            # Add slide number if requested
            if include_slide_numbers:
                slide_text_parts.append(f"=== SLIDE {slide_idx} ===")

            # Extract slide content
            slide_content = await self._extract_slide_content(
                slide, extract_table_content, extract_chart_titles
            )
            if slide_content:
                slide_text_parts.append(slide_content)

            # Extract speaker notes if requested
            if include_speaker_notes:
                notes_content = await self._extract_speaker_notes(slide)
                if notes_content:
                    slide_text_parts.append("--- Speaker Notes ---")
                    slide_text_parts.append(notes_content)

            if slide_text_parts:
                slide_contents.append("\n\n".join(slide_text_parts))

        return slide_separator.join(slide_contents)

    def _is_slide_hidden(self, slide: "Slide") -> bool:
        """Check if a slide is hidden."""
        try:
            # Check if slide has hidden property (this is a simplified check)
            # In practice, PowerPoint slide hiding is complex and may require
            # deeper XML parsing, but this covers basic cases
            slide_elem = slide._element
            return slide_elem.get("show") == "0"
        except (AttributeError, TypeError):
            return False

    async def _extract_slide_content(
        self, slide: "Slide", extract_table_content: bool, extract_chart_titles: bool
    ) -> str:
        """Extract text content from a single slide."""
        content_parts = []

        # Extract text from all shapes
        for shape in slide.shapes:
            try:
                # Handle text frames (text boxes, titles, content placeholders)
                if hasattr(shape, "text_frame") and shape.text_frame:
                    shape_text = await self._extract_text_frame_content(
                        shape.text_frame
                    )
                    if shape_text:
                        content_parts.append(shape_text)

                # Handle tables
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE and extract_table_content:
                    table_text = await self._extract_table_content(shape.table)
                    if table_text:
                        content_parts.append(table_text)

                # Handle charts (extract title if available)
                elif shape.shape_type == MSO_SHAPE_TYPE.CHART and extract_chart_titles:
                    chart_text = await self._extract_chart_content(shape)
                    if chart_text:
                        content_parts.append(chart_text)

                # Handle grouped shapes
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    group_text = await self._extract_group_content(
                        shape, extract_table_content
                    )
                    if group_text:
                        content_parts.append(group_text)

            except Exception as e:
                logger.warning(f"Failed to extract content from shape: {e}")
                continue

        return "\n\n".join(content_parts)

    async def _extract_text_frame_content(self, text_frame) -> str:
        """Extract text content from a text frame with paragraph structure."""
        if not text_frame.paragraphs:
            return ""

        paragraphs = []
        for paragraph in text_frame.paragraphs:
            if not paragraph.runs:
                continue

            # Extract text from runs
            paragraph_text = ""
            for run in paragraph.runs:
                if run.text:
                    paragraph_text += run.text

            if paragraph_text.strip():
                # Add bullet point formatting if applicable
                if paragraph.level > 0:
                    indent = "  " * paragraph.level
                    paragraphs.append(f"{indent}• {paragraph_text.strip()}")
                else:
                    paragraphs.append(paragraph_text.strip())

        return "\n".join(paragraphs)

    async def _extract_table_content(self, table) -> str:
        """Extract text content from a table."""
        if not table.rows:
            return ""

        table_rows = []
        for row in table.rows:
            row_cells = []
            for cell in row.cells:
                # Extract text from cell text frame
                cell_text = ""
                if hasattr(cell, "text_frame") and cell.text_frame:
                    cell_text = await self._extract_text_frame_content(cell.text_frame)
                elif hasattr(cell, "text"):
                    cell_text = cell.text

                row_cells.append(cell_text.strip().replace("\n", " "))

            if any(cell for cell in row_cells):  # Skip empty rows
                table_rows.append(" | ".join(row_cells))

        if table_rows:
            return f"[TABLE]\n{chr(10).join(table_rows)}\n[/TABLE]"

        return ""

    async def _extract_chart_content(self, shape) -> str:
        """Extract title and basic info from charts."""
        try:
            chart_parts = []

            # Try to get chart title
            if hasattr(shape, "chart") and shape.chart:
                chart = shape.chart
                if hasattr(chart, "chart_title") and chart.chart_title:
                    if (
                        hasattr(chart.chart_title, "text_frame")
                        and chart.chart_title.text_frame
                    ):
                        title_text = await self._extract_text_frame_content(
                            chart.chart_title.text_frame
                        )
                        if title_text:
                            chart_parts.append(f"[CHART] {title_text}")

            # If no specific title found, add generic chart indicator
            if not chart_parts:
                chart_parts.append("[CHART]")

            return "\n".join(chart_parts)

        except Exception as e:
            logger.warning(f"Failed to extract chart content: {e}")
            return "[CHART]"

    async def _extract_group_content(
        self, group_shape, extract_table_content: bool
    ) -> str:
        """Extract text content from grouped shapes."""
        group_parts = []

        try:
            for shape in group_shape.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    shape_text = await self._extract_text_frame_content(
                        shape.text_frame
                    )
                    if shape_text:
                        group_parts.append(shape_text)
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE and extract_table_content:
                    table_text = await self._extract_table_content(shape.table)
                    if table_text:
                        group_parts.append(table_text)
        except Exception as e:
            logger.warning(f"Failed to extract group content: {e}")

        return "\n".join(group_parts)

    async def _extract_speaker_notes(self, slide: "Slide") -> str:
        """Extract speaker notes from a slide."""
        try:
            if not slide.has_notes_slide:
                return ""

            notes_slide = slide.notes_slide
            if not notes_slide.notes_text_frame:
                return ""

            notes_text = await self._extract_text_frame_content(
                notes_slide.notes_text_frame
            )
            return notes_text.strip() if notes_text else ""

        except Exception as e:
            logger.warning(f"Failed to extract speaker notes: {e}")
            return ""

    async def _has_images(self, prs: "Presentation") -> bool:
        """Check if presentation contains images."""
        try:
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        return True
                    # Check for images in groups
                    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        for group_shape in shape.shapes:
                            if group_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                return True
            return False
        except Exception:
            return False

    async def _has_tables(self, prs: "Presentation") -> bool:
        """Check if presentation contains tables."""
        try:
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        return True
                    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        for group_shape in shape.shapes:
                            if group_shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                                return True
            return False
        except Exception:
            return False

    async def _has_charts(self, prs: "Presentation") -> bool:
        """Check if presentation contains charts."""
        try:
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        return True
                    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        for group_shape in shape.shapes:
                            if group_shape.shape_type == MSO_SHAPE_TYPE.CHART:
                                return True
            return False
        except Exception:
            return False

    async def _has_speaker_notes(self, prs: "Presentation") -> bool:
        """Check if presentation has speaker notes."""
        try:
            for slide in prs.slides:
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if (
                        notes_slide.notes_text_frame
                        and notes_slide.notes_text_frame.text.strip()
                    ):
                        return True
            return False
        except Exception:
            return False

    def get_parsing_options(self) -> dict[str, dict[str, Any]]:
        """Get available parsing options for PPTX files."""
        return {
            "include_speaker_notes": {
                "type": bool,
                "default": True,
                "description": "Include speaker notes in extracted text",
            },
            "include_slide_numbers": {
                "type": bool,
                "default": True,
                "description": "Include slide numbers in output",
            },
            "slide_separator": {
                "type": str,
                "default": "\n\n---\n\n",
                "description": "Separator text between slides",
            },
            "extract_table_content": {
                "type": bool,
                "default": True,
                "description": "Extract text content from tables",
            },
            "extract_chart_titles": {
                "type": bool,
                "default": True,
                "description": "Extract titles and labels from charts",
            },
            "include_hidden_slides": {
                "type": bool,
                "default": False,
                "description": "Include hidden slides in extraction",
            },
        }
